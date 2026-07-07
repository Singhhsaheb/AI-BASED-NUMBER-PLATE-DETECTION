from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Based Number Plate Detection (ANPR)"
    subtitle.text = "A Deep Learning & Computer Vision Approach\n\nB.Tech Project Presentation\nBy Sonu Kumar\nEnrollment No: 2522759050024"

    # 2. Problem Statement
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Traffic volume is increasing exponentially, making manual monitoring and toll collection highly inefficient."
    p = tf.add_paragraph()
    p.text = "Law enforcement needs automated systems to quickly identify stolen or speeding vehicles without manual intervention."
    p = tf.add_paragraph()
    p.text = "Solution: An Automatic Number Plate Recognition (ANPR) system that detects and reads license plates from images in real-time."

    # 3. Objectives
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Project Objectives"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "To develop a low-cost, highly accurate automated license plate reading system."
    p = tf.add_paragraph()
    p.text = "To apply modern Computer Vision techniques (filtering, edge detection) for robust plate localization."
    p = tf.add_paragraph()
    p.text = "To integrate Deep Learning (Optical Character Recognition) to extract text from the detected plates accurately."

    # 4. Technologies Used
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Technologies & Tools Used"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Python 3: Core programming language."
    p = tf.add_paragraph()
    p.text = "OpenCV (Open Source Computer Vision):"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Used for Grayscaling, Bilateral Filtering, and Canny Edge Detection."
    p2.level = 2
    p = tf.add_paragraph()
    p.text = "EasyOCR & PyTorch:"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Deep-learning framework used for extracting character strings from the cropped plate image."
    p2.level = 2

    # 5. Methodology / Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "System Architecture (How it Works)"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "1. Image Acquisition: Load the vehicle image."
    p = tf.add_paragraph()
    p.text = "2. Pre-processing: Convert to grayscale & apply bilateral filtering to reduce noise while keeping edges sharp."
    p = tf.add_paragraph()
    p.text = "3. Edge Detection: Apply Canny edge detector to find distinct boundaries."
    p = tf.add_paragraph()
    p.text = "4. Contour Extraction: Search for rectangular shapes (4 distinct points)."
    p = tf.add_paragraph()
    p.text = "5. Cropping: Mask the image and isolate the number plate."
    p = tf.add_paragraph()
    p.text = "6. OCR: Pass the isolated plate to EasyOCR for text extraction."

    # 6. Applications & Advantages
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Applications & Advantages"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Real-World Applications:"
    p = tf.add_paragraph()
    p.text = "Electronic Toll Collection (ETC) systems."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Smart Parking management systems."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Traffic law enforcement (speeding, red-light cameras)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Advantages:"
    p = tf.add_paragraph()
    p.text = "Highly scalable and does not require specialized hardware."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated process significantly reduces human error."
    p.level = 1

    # 7. Conclusion & Future Scope
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion & Future Scope"
    tf = body_shape.text_frame
    p = tf.add_paragraph()
    p.text = "Conclusion:"
    p = tf.add_paragraph()
    p.text = "Successfully implemented an AI-driven approach capable of identifying and reading license plates under various lighting conditions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Enhancements:"
    p = tf.add_paragraph()
    p.text = "Integration with live CCTV camera feeds."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Connecting the system to a database for automatic log generation."
    p.level = 1

    # 8. Thank You Slide
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Thank You!"
    
    # Add subtitle text box manually
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Open for Questions"
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = 2  # Center alignment

    # Save presentation
    prs.save("Number_Plate_Detection_Presentation.pptx")
    print("Updated B.Tech Presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
